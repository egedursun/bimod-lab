#  Copyright (c) 2024 BMD™ Autonomous Holdings. All rights reserved.
#
#  Project: Bimod.io™
#  File: function_utils.py
#  Last Modified: 2024-10-05 01:39:48
#  Author: Ege Dogan Dursun (Co-Founder & Chief Executive Officer / CEO @ BMD™ Autonomous Holdings)
#  Created: 2024-10-05 14:42:48
#
#  This software is proprietary and confidential. Unauthorized copying,
#  distribution, modification, or use of this software, whether for
#  commercial, academic, or any other purpose, is strictly prohibited
#  without the prior express written permission of BMD™ Autonomous
#  Holdings.
#
#   For permission inquiries, please contact: admin@Bimod.io.
#


import logging
from io import BytesIO

from docx import Document as ms_dx
from pptx import Presentation as ms_px
from openpyxl import load_workbook as ms_ex


logger = logging.getLogger(__name__)


def decode_stream__docx(file_bytes):
    f_stream = BytesIO(file_bytes)
    doc = ms_dx(f_stream)
    logger.info(f"Decoded {len(doc.paragraphs)} paragraphs.")
    return '\n'.join([paragraph.text for paragraph in doc.paragraphs])


def decode_stream__pptx(file_bytes):
    f_stream = BytesIO(file_bytes)
    pres = ms_px(f_stream)
    txt = []
    for sl in pres.slides:
        for shp in sl.shapes:
            if hasattr(shp, "text"):
                txt.append(shp.text)
    logger.info(f"Decoded {len(txt)} slides from {len(pres.slides)} slides.")
    return '\n'.join(txt)


def decode_stream__xlsx(file_bytes):
    f_stream = BytesIO(file_bytes)
    w_book = ms_ex(f_stream, data_only=True)
    txt = []
    for sheet in w_book.sheetnames:
        w_sheet = w_book[sheet]
        for row in w_sheet.iter_rows(values_only=True):
            txt.append('\t'.join([str(cell) if cell is not None else '' for cell in row]))
    logger.info(f"Decoded {len(txt)} rows from {len(w_book.sheetnames)} sheets.")
    return '\n'.join(txt)
